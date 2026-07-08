"""Generate professional 2-slide government HR PowerPoint for DocVerify."""

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "DocVerify_HR_Presentation.pptx"

# Government / enterprise palette
NAVY = RGBColor(15, 23, 42)
GREEN = RGBColor(22, 163, 74)
SLATE = RGBColor(71, 85, 105)
DARK = RGBColor(30, 41, 59)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(248, 250, 252)
BORDER = RGBColor(226, 232, 240)
ACCENT_BLUE = RGBColor(37, 99, 235)

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(0.95)
MARGIN = Inches(0.45)


def add_header_bar(slide, title: str, subtitle: str):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, HEADER_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(MARGIN, Inches(0.12), Inches(9.1), Inches(0.42))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub_box = slide.shapes.add_textbox(MARGIN, Inches(0.52), Inches(9.1), Inches(0.32))
    stf = sub_box.text_frame
    stf.text = subtitle
    sp = stf.paragraphs[0]
    sp.font.size = Pt(11)
    sp.font.color.rgb = RGBColor(203, 213, 225)


def add_section_box(slide, left, top, width, height, title: str, bullets: list[str]):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.bold = True
    p0.font.size = Pt(14)
    p0.font.color.rgb = GREEN
    p0.space_after = Pt(6)

    for line in bullets:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.space_after = Pt(4)
        p.level = 0


def add_arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = SLATE
    conn.line.width = Pt(1.5)


def add_flow_diagram(slide, left, top, width, height):
    label = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    label.text_frame.text = "Verification flow"
    label.text_frame.paragraphs[0].font.size = Pt(13)
    label.text_frame.paragraphs[0].font.bold = True
    label.text_frame.paragraphs[0].font.color.rgb = DARK

    steps = [
        ("1. Upload", "Photo / PDF"),
        ("2. OCR Read", "Extract text"),
        ("3. AI Check", "Validate & detect"),
        ("4. Score", "Confidence %"),
        ("5. Save", "Audit record"),
    ]
    box_w = Inches(1.05)
    box_h = Inches(0.78)
    gap = Inches(0.12)
    y = top + Inches(0.35)
    x = left
    boxes = []

    for i, (title, sub) in enumerate(steps):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = GREEN if i % 2 == 0 else ACCENT_BLUE
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(7)
        p2.font.color.rgb = SLATE
        p2.alignment = PP_ALIGN.CENTER
        boxes.append(box)
        x += box_w + gap

    for i in range(len(boxes) - 1):
        b1, b2 = boxes[i], boxes[i + 1]
        add_arrow(
            slide,
            b1.left + b1.width,
            b1.top + b1.height / 2,
            b2.left,
            b2.top + b2.height / 2,
        )


def add_document_table(slide, left, top, width, height):
    rows, cols = 7, 2
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    headers = ["Document type", "What HR receives"]
    data = [
        ("Aadhaar", "Number read, format checked, masked copy"),
        ("PAN", "PAN validated, name extracted"),
        ("Caste certificate", "Category and certificate details"),
        ("Experience letter", "Company, employee, dates"),
        ("Education", "Institute, degree, year"),
        ("Resume", "Contact, education, experience"),
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE

    for r, (doc, out) in enumerate(data, start=1):
        table.cell(r, 0).text = doc
        table.cell(r, 1).text = out
        for c in range(2):
            for p in table.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK


def add_architecture_row(slide, left, top, width):
    label = slide.shapes.add_textbox(left, top, width, Inches(0.25))
    label.text_frame.text = "System architecture"
    label.text_frame.paragraphs[0].font.size = Pt(13)
    label.text_frame.paragraphs[0].font.bold = True
    label.text_frame.paragraphs[0].font.color.rgb = DARK

    parts = [
        ("Netlify", "HR website", "Upload, dashboard, review"),
        ("HF Space", "Processing engine", "OCR + AI + validation"),
        ("Supabase", "Secure database", "Results and audit trail"),
    ]
    box_w = Inches(2.7)
    box_h = Inches(1.05)
    gap = Inches(0.35)
    y = top + Inches(0.32)
    x = left + Inches(0.15)
    boxes = []

    for title, role, desc in parts:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BORDER
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = role
        p2.font.size = Pt(10)
        p2.font.color.rgb = GREEN
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(8)
        p3.font.color.rgb = SLATE
        p3.alignment = PP_ALIGN.CENTER
        boxes.append(box)
        x += box_w + gap

    for i in range(len(boxes) - 1):
        b1, b2 = boxes[i], boxes[i + 1]
        add_arrow(
            slide,
            b1.left + b1.width,
            b1.top + b1.height / 2,
            b2.left,
            b2.top + b2.height / 2,
        )


def slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "DocVerify AI", "Government HR Document Verification")

    add_section_box(
        slide,
        MARGIN, Inches(1.15), Inches(4.2), Inches(2.35),
        "Purpose for HR",
        [
            "HR receives official papers: Aadhaar, PAN, caste, experience, education, resume.",
            "Manual checking is slow and risks errors or sharing private numbers.",
            "DocVerify uploads a photo/PDF and returns verified, review, or flagged.",
            "Sensitive data is masked; every result is saved for audit.",
        ],
    )

    add_flow_diagram(slide, Inches(4.85), Inches(1.15), Inches(4.7), Inches(1.2))

    add_document_table(slide, MARGIN, Inches(3.65), Inches(9.1), Inches(2.0))

    footer = slide.shapes.add_textbox(MARGIN, Inches(5.85), Inches(9.1), Inches(0.4))
    footer.text_frame.text = (
        "HR keeps final authority. Technology handles first-level reading; uncertain cases go to review queue."
    )
    footer.text_frame.paragraphs[0].font.size = Pt(10)
    footer.text_frame.paragraphs[0].font.italic = True
    footer.text_frame.paragraphs[0].font.color.rgb = SLATE


def slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "How DocVerify Works", "System overview for HR officers")

    add_architecture_row(slide, MARGIN, Inches(1.1), Inches(9.1))

    add_section_box(
        slide,
        MARGIN, Inches(2.55), Inches(4.4), Inches(2.5),
        "Process and outcomes",
        [
            "1. HR uploads document on the website.",
            "2. Processing engine reads text and runs AI + rule checks.",
            "3. Results stored in secure database with scores and flags.",
            "4. Final score: 70% AI analysis + 30% rule validation.",
            "Outcomes: Verified · Manual Review · Rejected.",
        ],
    )

    chart_data = CategoryChartData()
    chart_data.categories = ["OCR", "Fields", "Rules", "Image", "Final"]
    chart_data.add_series("Score (example)", (28, 24, 18, 12, 82))
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(5.05), Inches(2.55), Inches(4.5), Inches(2.5),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_title.text_frame.text = "Example trust score breakdown"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Points"
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "Check type"

    summary = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        MARGIN, Inches(5.25), Inches(9.1), Inches(0.55),
    )
    summary.fill.solid()
    summary.fill.fore_color.rgb = LIGHT
    summary.line.color.rgb = GREEN
    summary.line.width = Pt(1.5)
    stf = summary.text_frame
    stf.text = "Upload  →  Check  →  Mask private data  →  Save  →  HR reviews if needed"
    sp = stf.paragraphs[0]
    sp.font.bold = True
    sp.font.size = Pt(14)
    sp.font.color.rgb = NAVY
    sp.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide1(prs)
    slide2(prs)
    prs.save(OUT)
    print(f"Created: {OUT}")


if __name__ == "__main__":
    main()
